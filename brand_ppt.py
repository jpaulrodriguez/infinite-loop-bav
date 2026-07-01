"""
Brand PPT Generator — Genera el deck para cualquier marca del BRANDS array.
Usa brand_engine.py como fuente de verdad (anti-hallucination).
Mismo diseño que Makro_Brand_Intelligence_2026_v2.pptx.
"""
import io, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from brand_engine import analyze_brand

# ── LOGO VML ──────────────────────────────────────────────
_DIR = os.path.dirname(os.path.abspath(__file__))
VML_LOGO = os.path.join(_DIR, "vml_logo_white.png")

# ── PALETA ────────────────────────────────────────────────
BG        = RGBColor(0x07,0x0B,0x19)
BLACK     = RGBColor(0x00,0x00,0x00)
SURFACE   = RGBColor(0x12,0x1A,0x2F)
SURFACE_I = RGBColor(0x0A,0x10,0x1E)
SURFACE_3 = RGBColor(0x1A,0x24,0x3E)
WHITE     = RGBColor(0xFF,0xFF,0xFF)
OFFWHITE  = RGBColor(0xF0,0xF2,0xF8)
GRAY_1    = RGBColor(0xA3,0xA8,0xB8)
GRAY_2    = RGBColor(0x68,0x70,0x8A)
GRAY_3    = RGBColor(0x22,0x2C,0x44)
ORANGE    = RGBColor(0xFF,0x7A,0x00)
ORANGE_DIM= RGBColor(0x33,0x18,0x00)
GREEN     = RGBColor(0x00,0xD2,0x6A)
GREEN_DIM = RGBColor(0x00,0x2B,0x16)
RED       = RGBColor(0xFF,0x3B,0x30)
RED_DIM   = RGBColor(0x33,0x0B,0x09)

W = Inches(13.33)
H = Inches(7.5)

def _sem_color(sem: str):
    if sem == "verde":   return GREEN
    if sem == "rojo":    return RED
    if sem == "amarillo":return ORANGE
    return GRAY_1

def _sem_dim(sem: str):
    if sem == "verde":   return GREEN_DIM
    if sem == "rojo":    return RED_DIM
    if sem == "amarillo":return ORANGE_DIM
    return SURFACE_3

def _frente_color(f: str):
    f = f.lower()
    if "cx" in f: return RED
    if "co" in f: return ORANGE
    return GRAY_1

def dim_col(color, factor=8):
    h = str(color)
    r,g,b = int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
    return RGBColor(r//factor, g//factor, b//factor)

# ── PRIMITIVOS ────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(sl):
    bg = sl.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG

def rect(sl, x, y, w, h, fill, lc=None, lw=Pt(0.4)):
    shp = sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if lc: shp.line.color.rgb = lc; shp.line.width = lw
    else:  shp.line.fill.background()
    return shp

def oval_orb(sl, x, y, w, h, color, alpha=20):
    shp = sl.shapes.add_shape(9,Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    sf = shp._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if sf is not None:
        sc = sf.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if sc is not None:
            ae = etree.SubElement(sc,'{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            ae.set('val', str(alpha*1000))
    shp.line.fill.background()
    return shp

def pill(sl, x, y, w, h, fill, lc=None, lw=Pt(0.75)):
    shp = sl.shapes.add_shape(9,Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if lc: shp.line.color.rgb = lc; shp.line.width = lw
    else:  shp.line.fill.background()
    return shp

def txt(sl, text, x, y, w, h, size=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Montserrat"):
    tb = sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    run.font.name = font
    return tb

def eyebrow(sl, text, x, y, color=ORANGE):
    r,g,b = int(str(color)[0:2],16),int(str(color)[2:4],16),int(str(color)[4:6],16)
    pill(sl,x,y,len(text)*0.075+0.25,0.24,RGBColor(r//6,g//6,b//6),lc=color,lw=Pt(0.5))
    txt(sl,text.upper(),x+0.13,y+0.04,len(text)*0.1+0.2,0.2,size=7.5,bold=True,color=color,font="Montserrat")

def rule(sl, x, y, w, color=GRAY_3):
    rect(sl,x,y,w,0.01,color)

def card(sl, x, y, w, h, accent=None, g=0.08):
    rect(sl,x,y,w,h,SURFACE,lc=GRAY_3,lw=Pt(0.4))
    rect(sl,x+g,y+g,w-g*2,h-g*2,SURFACE_I)
    if accent: rect(sl,x,y,w,0.045,accent)
    return (x+g+0.12, y+g+0.18)

def stat_card(sl, x, y, w, h, label, value, sub, vc=WHITE, accent=None):
    tx,ty = card(sl,x,y,w,h,accent=accent)
    txt(sl,label.upper(),tx,ty,w-0.3,0.22,size=7,bold=True,color=GRAY_2)
    txt(sl,str(value),tx-0.05,ty+0.28,w-0.2,h*0.45,size=30,bold=True,color=vc)
    txt(sl,sub,tx,ty+0.28+h*0.42,w-0.2,0.25,size=8,color=GRAY_2)

def _truncate(text: str, n: int) -> str:
    """Trunca en límite de palabra y añade '…' si es necesario."""
    if len(text) <= n:
        return text
    cut = text[:n].rsplit(" ", 1)[0]
    return cut.rstrip(".,;:") + "…"

def add_logo(sl):
    if os.path.exists(VML_LOGO):
        lw = Inches(0.48); lh = Inches(0.48)
        lx = W - lw - Inches(0.22); ly = H - lh - Inches(0.15)
        sl.shapes.add_picture(VML_LOGO, lx, ly, lw, lh)

# ── CHARTS ────────────────────────────────────────────────
def chart_bars_h(sl, x, y, w, h, items, max_val=100):
    n = len(items); row_h = h/n; baw = w-2.2; lw = 2.0
    for i,(label,mv,cv,col) in enumerate(items):
        yi = y+i*row_h; rcy = yi+row_h/2
        txt(sl,label,x,rcy-0.18,lw-0.1,0.36,size=8.5,color=GRAY_1,align=PP_ALIGN.RIGHT)
        bx = x+lw+0.1
        rect(sl,bx,rcy-0.045,baw,0.09,SURFACE_3)
        cw = (cv/max_val)*baw if cv else 0
        rect(sl,bx,rcy-0.045,max(cw,0.02),0.045,GRAY_2)
        mw = (mv/max_val)*baw if mv else 0
        rect(sl,bx,rcy,max(mw,0.02),0.045,col)
        if cv: txt(sl,f"{int(cv)}",bx+cw+0.06,rcy-0.22,0.35,0.22,size=8,color=GRAY_2)
        if mv: txt(sl,f"{int(mv)}",bx+mw+0.06,rcy+0.04,0.35,0.22,size=8.5,bold=True,color=col)
    ly = y+h+0.12
    rect(sl,x+lw+0.1,ly+0.03,0.28,0.09,GRAY_2)
    txt(sl,"Líderes (Top 3)",x+lw+0.45,ly,1.8,0.22,size=7.5,color=GRAY_2)
    rect(sl,x+lw+2.4,ly+0.03,0.28,0.09,ORANGE)
    txt(sl,"Marca",x+lw+2.75,ly,1.0,0.22,size=7.5,bold=True,color=WHITE)

def chart_bars_v(sl, x, y, w, h, items, highlight=None):
    n = len(items); bw = (w-0.08*(n-1))/n; axh = h-0.55
    max_v = max((v for _,v in items),default=100)*1.15 or 100
    for g_v in [25,50,75,100]:
        if g_v <= max_v:
            gy = y+axh-(g_v/max_v)*axh
            rect(sl,x,gy,w,0.006,GRAY_3)
            txt(sl,str(g_v),x-0.38,gy-0.12,0.35,0.25,size=7,color=GRAY_2,align=PP_ALIGN.RIGHT)
    for i,(name,val) in enumerate(items):
        xi = x+i*(bw+0.08); bh = (val/max_v)*axh
        is_hl = name==highlight; col = ORANGE if is_hl else SURFACE_3
        txt_col = ORANGE if is_hl else GRAY_2
        rect(sl,xi+0.03,y+axh-bh+0.03,bw,bh,SURFACE)
        rect(sl,xi,y+axh-bh,bw,bh,col,lc=(ORANGE if is_hl else None),lw=Pt(0.4) if is_hl else None)
        txt(sl,str(int(val)),xi,y+axh-bh-0.28,bw,0.25,size=8,bold=is_hl,color=txt_col,align=PP_ALIGN.CENTER)
        txt(sl,name,xi-0.04,y+axh+0.08,bw+0.08,0.4,size=6.5,color=txt_col,align=PP_ALIGN.CENTER)
    rect(sl,x,y+axh,w,0.012,GRAY_1)

def chart_scatter(sl, x, y, w, h, brands, highlight=None):
    ax_x=x+0.55; ax_y=y; ax_w=w-0.55; ax_h=h-0.55
    rect(sl,ax_x,ax_y,ax_w/2,ax_h/2,RGBColor(0x10,0x18,0x2A))
    rect(sl,ax_x+ax_w/2,ax_y,ax_w/2,ax_h/2,RGBColor(0x0A,0x14,0x22))
    rect(sl,ax_x,ax_y+ax_h/2,ax_w/2,ax_h/2,RGBColor(0x14,0x0A,0x12))
    rect(sl,ax_x+ax_w/2,ax_y+ax_h/2,ax_w/2,ax_h/2,RGBColor(0x0A,0x12,0x10))
    for v in [25,50,75,100]:
        gx = ax_x+(v/100)*ax_w; gy = ax_y+ax_h-(v/100)*ax_h
        rect(sl,gx,ax_y,0.006,ax_h,GRAY_3)
        txt(sl,str(v),gx-0.15,ax_y+ax_h+0.06,0.3,0.22,size=7,color=GRAY_2,align=PP_ALIGN.CENTER)
        rect(sl,ax_x,gy,ax_w,0.006,GRAY_3)
        txt(sl,str(v),ax_x-0.42,gy-0.11,0.38,0.22,size=7,color=GRAY_2,align=PP_ALIGN.RIGHT)
    rect(sl,ax_x,ax_y+ax_h,ax_w,0.016,GRAY_1)
    rect(sl,ax_x,ax_y,0.016,ax_h,GRAY_1)
    qs = dict(size=7.5,color=GRAY_3,italic=True,font="Montserrat")
    txt(sl,"Zona crítica",ax_x+0.1,ax_y+ax_h*0.82,ax_w*0.45,0.25,**qs)
    txt(sl,"Marca sin experiencia",ax_x+ax_w*0.55,ax_y+ax_h*0.82,ax_w*0.4,0.25,**qs)
    txt(sl,"Potencial sin marca",ax_x+0.1,ax_y+0.05,ax_w*0.45,0.25,**qs)
    txt(sl,"Líderes del sector ↗",ax_x+ax_w*0.55,ax_y+0.05,ax_w*0.4,0.25,**qs)
    txt(sl,"Fuerza de Marca (BX)  →",ax_x+ax_w/2-1.3,ax_y+ax_h+0.28,2.6,0.28,size=8.5,color=GRAY_1,align=PP_ALIGN.CENTER)
    for bd in brands:
        nm = bd["name"]; bx_s = bd["bx"]; cx_s = bd["cx"]
        px = ax_x+(bx_s/100)*ax_w; py = ax_y+ax_h-(cx_s/100)*ax_h
        is_hl = nm==highlight; r = 0.18 if is_hl else 0.11
        if is_hl: oval_orb(sl,px-0.25,py-0.25,0.5,0.5,ORANGE,alpha=15)
        dot = sl.shapes.add_shape(9,Inches(px-r/2),Inches(py-r/2),Inches(r),Inches(r))
        dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE if is_hl else SURFACE_3
        if is_hl: dot.line.color.rgb = WHITE; dot.line.width = Pt(1.5)
        else:     dot.line.color.rgb = GRAY_2; dot.line.width = Pt(0.5)
        lbx = px+0.1 if bx_s < 80 else px-1.3
        txt(sl,nm,lbx,py-0.13,1.3,0.25,size=7 if not is_hl else 8.5,
            bold=is_hl,color=ORANGE if is_hl else GRAY_2)

# ── CONTACTOS POR ESPECIALIDAD ────────────────────────────
CONTACTS = {
    "cx": ("José Ocampo",        "jose.ocampo@vml.com",        "Experiencia del Cliente (CX)"),
    "bx": ("Paula Tejada",       "paula.tejada@vml.com",       "Fuerza de Marca (BX)"),
    "co": ("Jonathan Rodríguez", "jonathan.rodriguez@vml.com", "Conversión Comercial (CO)"),
}

def slide_specialty(prs, marca, eyebrow_txt, headline, ranking, analysis, col, caption):
    """Slide de una especialidad: ranking del sector (gráfica) + diagnóstico completo."""
    sl = blank(prs); fill_bg(sl)
    oval_orb(sl, 6, 0, 8, 8, col, alpha=6)
    eyebrow(sl, eyebrow_txt, 1.0, 0.35, color=col)
    txt(sl, headline, 1.0, 0.68, 11.5, 1.2, size=27, bold=True, color=WHITE)
    rule(sl, 1.0, 1.98, 11.33)
    items = [(d["name"], d["val"]) for d in ranking]
    if items:
        chart_bars_v(sl, 0.8, 2.15, 11.7, 3.35, items, highlight=marca)
        txt(sl, caption, 0.8, 5.55, 11.7, 0.28, size=7.5, color=GRAY_2, align=PP_ALIGN.CENTER)
    tx, ty = card(sl, 0.8, 5.85, 11.7, 1.45, accent=col, g=0.08)
    eyebrow(sl, "Diagnóstico", tx, ty + 0.02, color=col)
    txt(sl, _truncate(analysis, 440), tx, ty + 0.42, 11.2, 0.75, size=10, color=GRAY_1)
    add_logo(sl)

def slide_contact(prs, marca, primary_frente):
    """Slide de cierre: la persona adecuada según el frente prioritario de la marca."""
    sl = blank(prs); fill_bg(sl)
    oval_orb(sl, 4, 1, 9, 8, ORANGE, alpha=6)
    oval_orb(sl, -1.5, -1, 6, 6, RGBColor(0x00, 0x50, 0xFF), alpha=4)
    pf = primary_frente if primary_frente in CONTACTS else "cx"
    name, email, flabel = CONTACTS[pf]

    eyebrow(sl, "El siguiente paso", 1.0, 0.5, color=ORANGE)
    txt(sl, "Del diagnóstico\nal plan de acción.", 1.0, 0.95, 11, 1.7, size=42, bold=True, color=WHITE)
    rule(sl, 1.0, 3.05, 11.33)
    txt(sl, f"El frente prioritario para {marca} es {flabel}.\n"
            f"Para profundizar en un plan de acción enfocado, comunícate con:",
        1.0, 3.3, 11, 0.9, size=14, color=GRAY_1)

    # Contacto principal (el del tema a trabajar) — protagonista
    tx, ty = card(sl, 1.0, 4.5, 6.7, 2.05, accent=ORANGE, g=0.1)
    txt(sl, name, tx, ty + 0.12, 6.3, 0.7, size=32, bold=True, color=WHITE)
    txt(sl, email, tx, ty + 0.92, 6.3, 0.4, size=16, bold=True, color=ORANGE)
    txt(sl, flabel, tx, ty + 1.42, 6.3, 0.3, size=10, color=GRAY_2)

    # Otros frentes (referencia)
    txt(sl, "¿El tema es otro frente?", 8.15, 4.6, 4.4, 0.3, size=10, bold=True, color=GRAY_1)
    for i, f in enumerate([x for x in ["bx", "co", "cx"] if x != pf]):
        n2, e2, l2 = CONTACTS[f]
        yo = 5.05 + i * 0.72
        txt(sl, l2, 8.15, yo, 4.6, 0.25, size=9.5, bold=True, color=GRAY_1)
        txt(sl, f"{n2} · {e2}", 8.15, yo + 0.27, 4.6, 0.25, size=9, color=GRAY_2)

    txt(sl, "Capital Intangible · Brand Intelligence Loop · WPP Colombia 2026",
        1.0, 7.05, 9, 0.3, size=8, color=GRAY_2)
    add_logo(sl)

# ═══════════════════════════════════════════════════════════
# GENERADOR PRINCIPAL
# ═══════════════════════════════════════════════════════════

def generate_ppt(nombre: str) -> bytes:
    """
    Genera el PPT para una marca. Retorna bytes del archivo .pptx.
    Lanza ValueError si la marca no existe o tiene datos insuficientes.
    """
    # Análisis con anti-hallucination
    data = analyze_brand(nombre)
    b    = data["brand"]
    avg  = data["sector_avg"]
    gaps = data["gaps"]
    sem  = data["semaphores"]
    rnk  = data["rankings"]

    marca   = b["marca"]
    sector  = b["sector"]
    n_sec   = b["n_sector"]

    def s(f): return b.get(f)   # score (puede ser None)
    def a(f): return avg.get(f) # avg sector

    def gap_str(f):
        g = gaps.get(f)
        if g is None: return "vs líderes: sin dato"
        return f"{g:+.1f} vs líderes"

    def score_str(v, default="—"):
        return f"{v:.0f}" if v is not None else default

    # Colores por frente según semáforo
    bx_col = _sem_color(sem.get("bx","amarillo"))
    co_col = _sem_color(sem.get("co","amarillo"))
    cx_col = _sem_color(sem.get("cx","rojo"))

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── SLIDE 1: PORTADA ──────────────────────────────────
    sl = blank(prs); fill_bg(sl)
    oval_orb(sl,-1.5,-1.0,7.0,6.0,ORANGE,alpha=6)
    oval_orb(sl,7.5,3.5,7.0,6.0,RGBColor(0x00,0x50,0xFF),alpha=5)
    oval_orb(sl,3.0,5.0,5.0,4.0,GREEN,alpha=4)
    rect(sl,1.0,0.8,0.035,5.5,ORANGE)
    eyebrow(sl,f"Capital Intangible · {sector} · 2026",1.25,0.95)
    txt(sl,marca,1.25,1.45,8,2.1,size=72 if len(marca)<=10 else 52,bold=True,color=WHITE)
    txt(sl,"¿Cómo nos ve\nel mercado?",1.25,3.55,7,1.3,size=26,color=OFFWHITE)
    rule(sl,1.25,4.98,7.0,GRAY_3)
    txt(sl,f"Diagnóstico de Marca · Brand Intelligence Loop 2026",1.25,5.12,8,0.4,size=10,color=GRAY_2)
    for i,(label,v,col,_) in enumerate([
        (f"BX · Marca",        score_str(s("bx")),  bx_col, None),
        (f"CO · Ventas",       score_str(s("co")),  co_col, None),
        (f"CX · Experiencia",  score_str(s("cx")),  cx_col, None),
    ]):
        xi,yi = 10.0, 1.6+i*1.65
        card(sl,xi,yi,2.6,1.35,accent=col,g=0.07)
        txt(sl,label,xi+0.2,yi+0.22,2.2,0.25,size=7.5,color=GRAY_2,bold=True)
        txt(sl,v,xi+0.2,yi+0.5,2.0,0.72,size=38,bold=True,color=col)
        txt(sl,"/ 100",xi+0.95,yi+0.92,1.0,0.3,size=9,color=GRAY_2)
    txt(sl,f"Confidencial · Capital Intangible",1.25,7.08,6,0.28,size=8,color=GRAY_2)
    add_logo(sl)

    # ── SLIDE 2: IDEA DOMINANTE ───────────────────────────
    sl = blank(prs); fill_bg(sl)
    oval_orb(sl,5,1,8,6,RED,alpha=5)
    eyebrow(sl,"El diagnóstico en una línea",1.0,0.4)
    txt(sl,"“",0.8,0.7,1.0,1.4,size=80,color=ORANGE,bold=True)
    idea = data["dominant_idea"]
    # Cortar en líneas si es muy largo
    txt(sl,idea,1.3,0.9,9.5,2.8,size=22 if len(idea)<120 else 18,bold=False,color=OFFWHITE)
    txt(sl,"”",10.5,2.5,1.0,1.0,size=80,color=ORANGE,bold=True)
    rule(sl,4.15,1.0,11.3,GRAY_3)
    pillars = [
        ("BX\nMarca",    _truncate(data["bx_analysis"],235), f"BX · {score_str(s('bx'))} / 100", bx_col, 0.9),
        ("CO\nVentas",   _truncate(data["co_analysis"],235),  f"CO · {score_str(s('co'))} / 100", co_col, 4.7),
        ("CX\nExp.",     _truncate(data["cx_analysis"],235),  f"CX · {score_str(s('cx'))} / 100", cx_col, 8.5),
    ]
    for head,body_t,kpi,col,xi in pillars:
        tx,ty = card(sl,xi,4.3,3.65,2.88,accent=col,g=0.09)
        txt(sl,head,tx,ty+0.0,3.2,0.55,size=11.5,bold=True,color=WHITE)
        txt(sl,body_t,tx,ty+0.65,3.2,1.2,size=8.5,color=GRAY_1)
        pill(sl,tx,ty+1.95,1.4,0.28,dim_col(col,8),lc=col,lw=Pt(0.6))
        txt(sl,kpi,tx+0.12,ty+1.98,1.3,0.2,size=7.5,bold=True,color=col)
    add_logo(sl)

    # ── SLIDE 3: EL MODELO — descripciones personalizadas por marca ──
    sl = blank(prs); fill_bg(sl)
    oval_orb(sl,-2,2,8,6,RGBColor(0x00,0x50,0xFF),alpha=5)
    eyebrow(sl,f"Cómo leemos a {marca}",1.0,0.35)
    txt(sl,"Tres preguntas.\nTres respuestas reales.",1.0,0.68,7,1.3,size=32,bold=True,color=WHITE)
    rule(sl,1.0,2.12,11.33)
    fd = data["frente_descs"]
    frentes = [
        ("BX · Marca","¿Qué tan fuerte\nes la marca?",      fd["bx"], bx_col,s("bx"),sem.get("bx","amarillo")),
        ("CO · Ventas","¿Qué tanto\nvende?",                fd["co"], co_col,s("co"),sem.get("co","amarillo")),
        ("CX · Experiencia","¿Qué tan feliz\nqueda el cliente?", fd["cx"], cx_col,s("cx"),sem.get("cx","rojo")),
    ]
    sem_labels = {"verde":"Nivel alto","amarillo":"Nivel medio","rojo":"⚠ Alerta crítica","sin_dato":"Sin dato disponible"}
    for i,(frente,q,desc,col,score,semaf) in enumerate(frentes):
        xi = 0.7+i*4.22
        tx,ty = card(sl,xi,2.35,3.95,4.82,accent=col,g=0.1)
        eyebrow(sl,frente,xi+0.2,ty+0.05,color=col)
        txt(sl,q,xi+0.2,ty+0.42,3.4,0.9,size=14,bold=True,color=WHITE)
        rule(sl,xi+0.2,ty+1.42,3.4,col)
        txt(sl,_truncate(desc,290),xi+0.2,ty+1.58,3.4,1.5,size=8.5,color=GRAY_1)
        dot = sl.shapes.add_shape(9,Inches(xi+1.2),Inches(ty+3.15),Inches(1.55),Inches(1.15))
        dot.fill.solid(); dot.fill.fore_color.rgb = SURFACE
        dot.line.color.rgb = col; dot.line.width = Pt(1.5)
        sc = score_str(score,"—")
        txt(sl,sc,xi+1.25,ty+3.18,1.45,0.75,size=32,bold=True,color=col,align=PP_ALIGN.CENTER)
        txt(sl,sem_labels.get(semaf,"—"),xi+0.1,ty+4.38,3.6,0.28,size=7.5,color=col,bold=True,align=PP_ALIGN.CENTER)
    add_logo(sl)

    # ── SLIDE 4: MARCA VS SECTOR — BARRAS H ───────────────
    sl = blank(prs); fill_bg(sl)
    eyebrow(sl,f"{marca} vs. el sector · Todos los indicadores",1.0,0.35)
    txt(sl,"Dónde estamos y\ndónde está la competencia.",1.0,0.68,11,1.2,size=28,bold=True,color=WHITE)
    rule(sl,1.0,2.0,11.33)

    def _row(label, field, col):
        mv = s(field); cv = a(field)
        return (label, mv if mv is not None else 0, cv if cv is not None else 0, col)

    chart_items = [
        _row("Fuerza de Marca\n(Brand Asset)",  "brand_asset", GRAY_1),
        _row("Presencia Digital\n(BAV Pulse)",   "bav_pulse",   GREEN),
        _row("Volumen de Ventas\n(Commerce)",     "commerce",    co_col),
        _row("Experiencia en Punto\n(CX KPI)",    "cx_kpi",      cx_col),
        _row("Clientes que Vuelven\n(Loyalty)",   "loyalty",     cx_col),
    ]
    chart_bars_h(sl,0.8,2.2,11.7,4.35,chart_items)
    tx,ty = card(sl,0.8,6.8,11.7,0.57,g=0.06)
    txt(sl,_truncate(data["comparative_insight"],280),tx,ty+0.04,11.1,0.4,size=9,color=GRAY_1)
    add_logo(sl)

    # ── SLIDES 5-7: LAS TRES ESPECIALIDADES (diagnóstico completo + ranking del sector) ──
    def _rank_headline(nombre_frente, fkey):
        rk = rnk.get(fkey, {})
        pos = rk.get("pos"); tot = rk.get("total", n_sec)
        return (f"{nombre_frente}: puesto {pos} de {tot} del sector."
                if pos else f"{nombre_frente}: diagnóstico del frente.")

    slide_specialty(prs, marca, "BX · Fuerza de Marca",
        _rank_headline("Fuerza de marca", "bx"),
        data["bx_ranking"], data["bx_analysis"], bx_col,
        f"Marcas de {sector} ordenadas por Fuerza de Marca (BX) — de menor a mayor")
    slide_specialty(prs, marca, "CO · Conversión Comercial",
        _rank_headline("Conversión", "co"),
        data["co_ranking"], data["co_analysis"], co_col,
        f"Marcas de {sector} ordenadas por Commerce Score (CO) — de menor a mayor")
    slide_specialty(prs, marca, "CX · Experiencia del Cliente",
        _rank_headline("Experiencia", "cx"),
        data["cx_ranking"], data["cx_analysis"], cx_col,
        f"Marcas de {sector} ordenadas por Experiencia (CX) — de menor a mayor")

    # ── SLIDE 8: MAYOR FORTALEZA RELATIVA ────────────────
    sl = blank(prs); fill_bg(sl)
    bf   = data["strength_field"]
    bg   = data["strength_gap"]
    field_names = {"bx":"Fuerza de Marca","co":"Commerce Score","cx":"Experiencia del Cliente",
                   "brand_asset":"Activo de Marca","bav_pulse":"Presencia Digital"}
    best_name  = field_names.get(bf,"Indicador")
    best_score = score_str(s(bf))
    best_sector= score_str(a(bf))
    best_col   = GREEN if bg > 5 else (ORANGE if bg >= 0 else RED)

    oval_orb(sl,5,-1,9,8,best_col,alpha=6)
    heading_sl5 = "Lo que sí funciona ✓" if bg >= 0 else "El frente más cercano al sector"
    eyebrow(sl,heading_sl5,1.0,0.35,color=best_col)
    txt(sl,f"{best_name}:\n{'ventaja relativa' if bg > 0 else 'menor brecha'}.",1.0,0.68,9,1.2,size=30,bold=True,color=WHITE)
    rule(sl,1.0,1.98,11.33)

    for xi,val,label,col,sub in [
        (1.2, best_score, marca, best_col, f"{best_name} · {marca} · BAV 2026"),
        (7.0, best_sector, "Líderes\ndel Sector", GRAY_1, f"{best_name} · Líderes (Top 3)"),
    ]:
        tx,ty = card(sl,xi,2.2,4.5,3.5,accent=col)
        txt(sl,val,tx+0.1,ty+0.05,3.8,1.8,size=88,bold=True,color=col,align=PP_ALIGN.CENTER)
        txt(sl,label,tx+0.1,ty+1.95,3.8,0.55,size=14,bold=True,color=col,align=PP_ALIGN.CENTER)
        txt(sl,sub,tx+0.1,ty+2.55,3.8,0.35,size=8,color=GRAY_2,align=PP_ALIGN.CENTER)

    gap_disp = f"+{bg:.0f}" if bg >= 0 else f"{bg:.0f}"
    txt(sl,gap_disp,5.9,3.15,1.0,0.8,size=28,bold=True,color=best_col,align=PP_ALIGN.CENTER)
    txt(sl,"puntos\narriba" if bg>=0 else "puntos\nabajo",5.78,3.9,1.25,0.55,size=8,color=best_col,align=PP_ALIGN.CENTER)
    rect(sl,6.1,4.55,0.55,0.015,best_col)

    rule(sl,1.0,5.95,11.33)
    txt(sl,"¿Qué significa?",1.0,6.08,5,0.3,size=10,bold=True,color=WHITE)
    txt(sl,_truncate(data["strength_insight"],310),1.0,6.42,11.3,0.9,size=9.5,color=GRAY_1)
    add_logo(sl)

    # ── SLIDE 9: MAPA DE POSICIONAMIENTO ──────────────────
    sl = blank(prs); fill_bg(sl)
    eyebrow(sl,f"Mapa de posicionamiento · {sector}",1.0,0.35)
    txt(sl,"Fuerza de Marca\nvs. Experiencia del Cliente.",1.0,0.68,9,1.2,size=28,bold=True,color=WHITE)
    rule(sl,1.0,1.98,11.33)
    chart_scatter(sl,0.6,2.15,11.5,4.85,data["scatter"],highlight=marca)
    tx,ty = card(sl,0.6,7.1,11.5,0.3,g=0.05)
    # Determinar cuadrante actual
    bx_v = s("bx") or 0; cx_v = s("cx") or 0
    if bx_v >= 50 and cx_v >= 50:   cuad = "cuadrante de liderazgo (BX alto + CX alto). El reto es mantenerse."
    elif bx_v >= 50 and cx_v < 50:  cuad = "cuadrante de marca sin experiencia. Hay reconocimiento pero la vivencia no lo sostiene."
    elif bx_v < 50 and cx_v >= 50:  cuad = "cuadrante de potencial sin marca. La experiencia es buena pero la marca no es suficientemente conocida."
    else:                            cuad = "cuadrante crítico (BX bajo + CX bajo). El trabajo es construir desde los fundamentos."
    txt(sl,f"{marca} está en el {cuad}",tx,ty+0.04,11.0,0.22,size=8.5,color=GRAY_1)
    add_logo(sl)

    # ── SLIDE 9: LOS 3 MOVIMIENTOS ────────────────────────
    sl = blank(prs); fill_bg(sl)
    oval_orb(sl,3,2,7,7,ORANGE,alpha=4)
    eyebrow(sl,"Orientación estratégica · ¿Qué hacemos ahora?",1.0,0.35)
    txt(sl,"Tres movimientos.\nEn orden de impacto.",1.0,0.68,9,1.2,size=28,bold=True,color=WHITE)
    rule(sl,1.0,1.98,11.33)
    timing_labels = ["90 días","3 a 9 meses","9 a 18 meses"]
    move_colors = [RED,ORANGE,GREEN]
    for idx,mov in enumerate(data["movements"][:3]):
        xi = 0.7+idx*4.22
        col = move_colors[idx]
        tx,ty = card(sl,xi,2.18,3.95,5.1,accent=col,g=0.1)
        txt(sl,mov["num"],xi+0.18,ty-0.05,1.2,0.95,size=40,bold=True,color=col)
        pill(sl,xi+0.18,ty+0.92,1.2,0.4,dim_col(col,8),lc=col,lw=Pt(0.5))
        txt(sl,mov["timing"],xi+0.25,ty+0.97,1.1,0.32,size=6.5,bold=True,color=col,align=PP_ALIGN.CENTER)
        rule(sl,xi+0.18,ty+1.45,3.4,col)
        txt(sl,mov["title"],xi+0.18,ty+1.6,3.4,0.55,size=10.5,bold=True,color=WHITE)
        txt(sl,_truncate(mov["body"],255),xi+0.18,ty+2.2,3.4,1.5,size=8.5,color=GRAY_1)
        rule(sl,xi+0.18,ty+3.75,3.4,GRAY_3)
        txt(sl,mov["meta"],xi+0.18,ty+3.9,3.4,0.75,size=8,color=col,bold=True)
    add_logo(sl)

    # ── SLIDE 11: CONTACTO — la persona adecuada para el tema principal ──
    slide_contact(prs, marca, data["primary_frente"])

    # ── SERIALIZAR ────────────────────────────────────────
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
